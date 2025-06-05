import os
import numpy as np
import PyPDF2
from sentence_transformers import SentenceTransformer
from docx import Document as DocxDocument
import re

# Add imports for Excel and PowerPoint
import openpyxl
from pptx import Presentation

def extract_sources_list(source_docs):
    seen = set()
    sources = []
    for doc in source_docs:
        file_name = doc.metadata.get("file_name")
        if file_name and file_name not in seen:
            seen.add(file_name)
            sources.append(file_name)
    return sources

def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    print(f"[DEBUG] Processing file: {filepath}, type: {ext}")
    if ext == '.txt':
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
            print(f"[DEBUG] TXT extracted length: {len(text)}")
            return text
    elif ext == '.docx':
        doc = DocxDocument(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    tables.append(" | ".join(row_text))
        all_text = paragraphs + tables
        text = '\n'.join(all_text)
        print(f"[DEBUG] DOCX extracted length: {len(text)}")
        return text
    elif ext == '.pdf':
        text = ''
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ''
        print(f"[DEBUG] PDF extracted length: {len(text)}")
        return text
    elif ext in ['.xlsx', '.xls']:
        try:
            wb = openpyxl.load_workbook(filepath, data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text.append(' | '.join(row_text))
            result = '\n'.join(text)
            print(f"[DEBUG] XLSX/XLS extracted length: {len(result)}")
            return result
        except Exception as e:
            print(f"[ERROR] Excel extraction failed: {e}")
            return ''
    elif ext == '.pptx':
        try:
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = shape.text.strip()
                        if t:
                            text.append(t)
            result = '\n'.join(text)
            print(f"[DEBUG] PPTX extracted length: {len(result)}")
            return result
        except Exception as e:
            print(f"[ERROR] PPTX extraction: {e}")
            return ''
    else:
        print(f"[DEBUG] Unsupported file type: {ext}")
        return ''

def cosine_similarity(a, b):
    if not a.any() or not b.any():
        return 0.0
    return float(np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b)))

def normalize_text(text):
    text = text.lower().strip()
    text = re.sub(r'\s+', ' ', text)
    return text

def find_similar_files(uploaded_text, folder, threshold=0.7):
    model = SentenceTransformer('all-MiniLM-L12-v2')
    uploaded_text_norm = normalize_text(uploaded_text)
    uploaded_emb = model.encode([uploaded_text_norm], convert_to_numpy=True)[0]
    similar = []
    for fname in os.listdir(folder):
        fpath = os.path.join(folder, fname)
        if not os.path.isfile(fpath) or not fname.lower().endswith((".docx", ".pdf", ".txt", ".xlsx", ".pptx")):
            continue
        try:
            text = extract_text_from_file(fpath)
            text_norm = normalize_text(text)
            if not text_norm.strip():
                continue
            emb = model.encode([text_norm], convert_to_numpy=True)[0]
            sim = cosine_similarity(uploaded_emb, emb)
            if sim >= threshold:
                similar.append({
                    'file': fname,
                    'similarity': round(sim * 100, 2)
                })
        except Exception as e:
            print(f"[ERROR] Processing {fname}: {e}")
            continue
    similar.sort(key=lambda x: x['similarity'], reverse=True)
    exact_matches = [f for f in similar if f['similarity'] == 100.0]
    if exact_matches:
        return exact_matches
    return similar[:3]
